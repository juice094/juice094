from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

RUST = RGBColor(0xCE, 0x42, 0x2B)
DARK = RGBColor(0x11, 0x11, 0x14)
DARK_CARD = RGBColor(0x1E, 0x1E, 0x24)
GRAY = RGBColor(0x9C, 0xA3, 0xAF)
LIGHT_GRAY = RGBColor(0xD1, 0xD5, 0xDB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLUE = RGBColor(0x3B, 0x82, 0xF6)
GREEN = RGBColor(0x10, 0xB9, 0x81)

def set_bg(slide, color=DARK):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def text(slide, left, top, width, height, content, size=18, bold=False,
         color=WHITE, align=PP_ALIGN.LEFT, font="Microsoft YaHei", wrap=True):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.text = content
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font
    p.alignment = align
    return box

def line(slide, left, top, width, color=RUST, thickness=3):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top),
                                    Inches(width), Inches(thickness/72))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def card(slide, left, top, width, height, bg=DARK_CARD, radius=0.08):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top),
                                    Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg
    shape.line.fill.background()
    try:
        shape.adjustments[0] = radius
    except:
        pass
    return shape

def badge(slide, left, top, content, bg=RUST):
    w = min(max(len(content) * 0.15 + 0.3, 1.1), 3.0)
    card(slide, left, top, w, 0.38, bg)
    text(slide, left, top + 0.04, w, 0.3, content, size=12, bold=True,
         color=WHITE, align=PP_ALIGN.CENTER)
    return w

def big_num(slide, left, top, num, label, sub=""):
    text(slide, left, top, 2.8, 0.8, num, size=42, bold=True, color=RUST)
    text(slide, left, top + 0.75, 2.8, 0.35, label, size=14, color=LIGHT_GRAY)
    if sub:
        text(slide, left, top + 1.05, 2.8, 0.3, sub, size=11, color=GRAY)

def bullets(slide, left, top, width, height, items, size=15, color=LIGHT_GRAY, spacing=0.7):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Microsoft YaHei"
        p.space_after = Pt(spacing * 12)
    return box

# ===== 第1页：封面 =====
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)

bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                          Inches(0.12), Inches(7.5))
bar.fill.solid(); bar.fill.fore_color.rgb = RUST; bar.line.fill.background()

deco = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(10.8), Inches(-0.5),
                           Inches(3.0), Inches(2.8))
deco.fill.solid(); deco.fill.fore_color.rgb = RUST; deco.line.fill.background()
deco.rotation = -12

text(s, 0.8, 2.0, 10, 1.2, "周景潇", size=64, bold=True)
badge(s, 0.8, 3.35, "AI Agent 开发")
badge(s, 3.0, 3.35, "后端系统开发", DARK_CARD)
badge(s, 5.2, 3.35, "智能体平台研发", DARK_CARD)

text(s, 0.8, 4.15, 11, 0.7,
     "大三学生，用 Rust 独立做了三个系统级项目，专注本地优先和零依赖部署",
     size=22, color=GRAY)

line(s, 0.8, 5.7, 8, GRAY, 1)
text(s, 0.8, 5.9, 11, 0.5,
     "甘肃农业大学 · 数据科学与大数据技术 · 本科 2027 届",
     size=15, color=GRAY)
text(s, 0.8, 6.4, 11, 0.5,
     "13626566112  ·  2241470466@qq.com  ·  github.com/juice094",
     size=14, color=LIGHT_GRAY)

# ===== 第2页：我是谁 =====
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)

text(s, 0.8, 0.5, 10, 0.8, "我是谁", size=36, bold=True)
line(s, 0.8, 1.25, 1.5, RUST, 3)

text(s, 0.8, 1.8, 7.5, 2.6,
     "我主要用 Rust 写系统级软件。过去一年独立做了三个项目：一个跨设备保持一致的本地 AI Agent、"
     "一个不需要云盘的 P2P 文件同步工具、一个把代码和笔记整理成 AI 能直接用的上下文的知识库。\n\n"
     "我的工程习惯是：能本地跑就不依赖外部服务，能单二进制分发就不堆容器，生产代码里不写 unwrap。",
     size=18, color=LIGHT_GRAY)

big_num(s, 9.3, 1.8, "3", "独立项目")
big_num(s, 9.3, 3.5, "40+", "workspace crates")
big_num(s, 9.3, 5.2, "0", "生产 panic")

# ===== 第3页：我做项目遵循的思路 =====
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)

text(s, 0.8, 0.5, 10, 0.8, "我做项目遵循的思路", size=36, bold=True)
line(s, 0.8, 1.25, 3, RUST, 3)

lines = [
    ("本地优先", "用户数据留在自己设备上，推理和记忆也尽量本地完成，减少对外部 API 和数据库的依赖。",
     "Clarity 用 SQLite 存记忆，devbase 用 Tantivy 做本地搜索"),
    ("零外部依赖", "一个可执行文件就能跑，不强迫用户装 Python、Node.js 或向量数据库。",
     "devbase 部署从 4 步降到 1 步，Clarity cargo install 即可运行"),
    ("纯 Rust 全栈", "从协议编解码到 UI 到移动端 FFI 都用 Rust，统一工具链、统一性能基线。",
     "BEP 协议、MCP 协议、Agent loop 全部自己实现"),
]

for i, (title, desc, example) in enumerate(lines):
    x = 0.7 + i * 4.2
    card(s, x, 1.7, 3.9, 4.6)
    text(s, x + 0.2, 1.9, 3.5, 0.6, title, size=22, bold=True)
    text(s, x + 0.2, 2.7, 3.5, 1.8, desc, size=15, color=LIGHT_GRAY)
    text(s, x + 0.2, 4.7, 3.5, 1.4, f"例子：{example}", size=13, color=GRAY)

# ===== 第4页：Clarity =====
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)

text(s, 0.8, 0.5, 9, 0.8, "Clarity", size=40, bold=True)
text(s, 0.8, 1.15, 10, 0.5,
     "本地 AI Agent 运行时：让同一个 AI 助手在手机、电脑、服务器上记住你是谁", size=18, color=GRAY)
badge(s, 10.5, 0.6, "核心项目")

# 问题
text(s, 0.8, 1.9, 11.7, 0.5, "我解决的问题", size=16, bold=True, color=RUST)
text(s, 0.8, 2.35, 11.7, 0.8,
     "现有的 AI 助手换个设备启动就记不住你，角色设定也会断。我想让同一个 AI 身份在不同设备上保持一致。",
     size=16, color=LIGHT_GRAY)

# 做法
bullets(s, 0.8, 3.2, 6.5, 3.2,
        ["把 AI 身份拆成四层：身份认同、人格配置、记忆图谱、设备分身",
         "每个设备分身本地自治，断网也能继续工作",
         "重连后只同步脱敏后的记忆增量，原始隐私数据不上云",
         "同一个身份可以跑不同模型：云端大模型、本地小模型、移动端模型"],
        size=14, color=LIGHT_GRAY, spacing=0.8)

# 数据
big_num(s, 7.8, 3.2, "23", "workspace crates")
big_num(s, 7.8, 4.8, "152K", "行 Rust 代码", "LOC")
big_num(s, 10.5, 3.2, "1,243", "测试")
big_num(s, 10.5, 4.8, "0", "生产 panic")

# ===== 第5页：syncthing-rust =====
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)

text(s, 0.8, 0.5, 10, 0.8, "syncthing-rust", size=40, bold=True)
text(s, 0.8, 1.15, 10, 0.5,
     "用 Rust 重写 Syncthing：让多台设备直接同步文件，不需要云盘", size=18, color=GRAY)

text(s, 0.8, 1.9, 11.7, 0.5, "我解决的问题", size=16, bold=True, color=RUST)
text(s, 0.8, 2.35, 11.7, 0.8,
     "我想验证自己能不能从零实现一个复杂的分布式协议，并且和官方 Go 版真正互通。",
     size=16, color=LIGHT_GRAY)

bullets(s, 0.8, 3.2, 11.7, 2.2,
        ["用 prost 实现 BEP 协议 Protobuf 编解码，用 rustls 做 TLS 1.3 传输",
         "实现 UDP 本地广播、STUN、UPnP、Relay 四种 NAT 穿透方式",
         "块级增量同步：只传变更的块，不是整个文件",
         "和官方 Go 版 Syncthing 做了 wire_compat 集成测试，验证线路格式兼容"],
        size=14, color=LIGHT_GRAY, spacing=0.7)

metrics = [("8", "crates"), ("5", "binaries"), ("59K", "LOC"), ("~12MB", "单二进制"), ("170MB", "稳定内存")]
for i, (num, label) in enumerate(metrics):
    x = 0.8 + i * 2.45
    text(s, x, 5.8, 2.3, 0.6, num, size=24, bold=True, color=RUST, align=PP_ALIGN.CENTER)
    text(s, x, 6.35, 2.3, 0.4, label, size=12, color=GRAY, align=PP_ALIGN.CENTER)

# ===== 第6页：devbase =====
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)

text(s, 0.8, 0.5, 10, 0.8, "devbase", size=40, bold=True)
text(s, 0.8, 1.15, 10, 0.5,
     "开发者知识库：把 Git 仓库、笔记、脚本整理成 AI 能直接用的上下文", size=18, color=GRAY)

text(s, 0.8, 1.9, 11.7, 0.5, "我解决的问题", size=16, bold=True, color=RUST)
text(s, 0.8, 2.35, 11.7, 0.8,
     "AI 客户端每次只能拿到当前文件内容，看不到项目历史、笔记、工作流。我想把整个工作空间变成 AI 能推理的结构化上下文。",
     size=16, color=LIGHT_GRAY)

# 三层流程
steps = [
    ("1", "解析", "tree-sitter 解析多语言 AST，提取代码符号和关系"),
    ("2", "索引", "Tantivy 做全文搜索，SQLite 存 embedding 和图关系"),
    ("3", "编排", "YAML DAG 工作流把 Skill、子工作流、并行、判断串起来"),
]
for i, (num, title, desc) in enumerate(steps):
    x = 0.7 + i * 4.2
    card(s, x, 3.3, 3.9, 2.4)
    text(s, x + 0.2, 3.45, 3.5, 0.6, f"{num}. {title}", size=20, bold=True)
    text(s, x + 0.2, 4.15, 3.5, 1.3, desc, size=14, color=GRAY)

big_num(s, 0.8, 6.1, "71", "MCP 工具")
big_num(s, 4.0, 6.1, "56K", "LOC")
big_num(s, 7.2, 6.1, "8.7MB", "单二进制")

# ===== 第7页：三个项目的关系 =====
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)

text(s, 0.8, 0.5, 10, 0.8, "三个项目是怎么连起来的", size=36, bold=True)
line(s, 0.8, 1.25, 3.5, RUST, 3)

nodes = [
    ("Clarity", "AI Agent 运行时\n管理人格、记忆、工具", 0.8, 2.3, RUST),
    ("syncthing-rust", "P2P 同步协议\n同步文件和数据", 4.8, 2.3, BLUE),
    ("devbase", "开发者知识库\n整理代码、笔记、工作流", 8.8, 2.3, GREEN),
]
for title, desc, x, y, color in nodes:
    card(s, x, y, 3.6, 1.6, color)
    text(s, x + 0.15, y + 0.15, 3.3, 0.45, title, size=18, bold=True)
    text(s, x + 0.15, y + 0.65, 3.3, 0.8, desc, size=13, color=WHITE)

line(s, 4.45, 3.05, 0.4, GRAY, 2)
line(s, 8.45, 3.05, 0.4, GRAY, 2)

relations = [
    "devbase 用 syncthing-rust 做工作空间的 P2P 同步",
    "Clarity 和 devbase 都用 MCP 协议暴露工具，AI 客户端调用方式一致",
    "三个项目共享同一套工程基线：SQLite 内嵌、零外部依赖、生产代码零 panic",
]
bullets(s, 1.5, 4.4, 10.3, 1.8, relations, size=16, color=LIGHT_GRAY, spacing=1.0)

# ===== 第8页：研究经历 =====
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)

text(s, 0.8, 0.5, 10, 0.8, "研究经历", size=36, bold=True)
line(s, 0.8, 1.25, 1.5, RUST, 3)

text(s, 0.8, 1.7, 11, 0.6,
     "检索增强生成中输出结构的实证研究", size=24, bold=True)
text(s, 0.8, 2.3, 8, 0.4, "github.com/juice094/acr-select", size=14, color=BLUE)

text(s, 0.8, 2.9, 8, 1.6,
     "我在做 RAG 系统时发现，输出格式（JSON、Markdown、自然语言）会明显影响生成质量。"
     "但现有研究大多只关注检索准确率，很少系统性地验证这个交互效应。",
     size=17, color=LIGHT_GRAY)

text(s, 0.8, 4.5, 8, 1.0,
     "我设计了一套多因子受控实验，固定检索内容，只改变输出格式和检索深度，观察模型输出质量的变化。",
     size=17, color=LIGHT_GRAY)

data = [("40", "实验条件"), ("3000+", "样本"), ("4", "模型架构"), ("论文", "撰写中")]
for i, (num, label) in enumerate(data):
    x = 0.8 + i * 3.1
    card(s, x, 5.5, 2.8, 1.5)
    text(s, x, 5.7, 2.8, 0.7, num, size=36, bold=True, color=RUST, align=PP_ALIGN.CENTER)
    text(s, x, 6.45, 2.8, 0.35, label, size=14, color=GRAY, align=PP_ALIGN.CENTER)

# ===== 第9页：其他 =====
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)

text(s, 0.8, 0.5, 10, 0.8, "认证与其他", size=36, bold=True)
line(s, 0.8, 1.25, 2, RUST, 3)

items = [
    "AI 应用工程师（中级）— 工业和信息化部认证",
    "日语专业四级（CJT-4）· 大学英语四级（CET-4）",
    "到岗时间：2026.07；可实习 3–6 个月，每周 5 天",
    "有公网部署经验：Clarity 前端在 Vercel，syncthing-rust 节点在日本 VPS",
]
bullets(s, 0.8, 1.8, 11.7, 3.5, items, size=18, color=LIGHT_GRAY, spacing=1.0)

# ===== 第10页：结尾 =====
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)

text(s, 0, 2.2, 13.333, 1.0, "期待加入你的团队", size=36, bold=True,
     color=WHITE, align=PP_ALIGN.CENTER)
text(s, 0, 3.1, 13.333, 0.8, "一起做扎实的 AI Agent 基础设施",
     size=26, bold=True, color=RUST, align=PP_ALIGN.CENTER)

line(s, 5.0, 4.15, 3.333, GRAY, 1)

text(s, 0, 4.6, 13.333, 0.5,
     "到岗时间：2026.07  ·  可实习 3–6 个月，每周 5 天",
     size=16, color=GRAY, align=PP_ALIGN.CENTER)
text(s, 0, 5.2, 13.333, 0.5,
     "13626566112  ·  2241470466@qq.com  ·  github.com/juice094",
     size=15, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

out = r'C:\Users\22414\dev\juice094\resume\周景潇-求职综合展示-v4.pptx'
prs.save(out)
print(f"PPT 已生成：{out}")
